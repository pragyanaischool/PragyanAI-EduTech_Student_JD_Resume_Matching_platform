import io
import zipfile
import pytest
from unittest.mock import patch, MagicMock
import fitz  # PyMuPDF
import docx
import pptx

from core.parsers import (
    parse_pdf,
    parse_docx,
    parse_pptx,
    parse_url,
    parse_zip,
    extract_text_auto
)


# ---------------------------------------------------------------------------
# Test Fixtures: Generating Valid In-Memory Binary Files
# ---------------------------------------------------------------------------

@pytest.fixture
def sample_pdf_bytes() -> bytes:
    """Generates a valid single-page PDF in-memory using PyMuPDF."""
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 72), "Aarav Sharma - Generative AI Specialist\nProficient in Python, LangGraph, and Groq.")
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


@pytest.fixture
def sample_docx_bytes() -> bytes:
    """Generates a valid DOCX file in-memory with paragraphs and tables."""
    doc = docx.Document()
    doc.add_heading("Lead Systems Architect", level=1)
    doc.add_paragraph("Specializing in ChromaDB and FAISS local vector stores.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Skill"
    table.cell(0, 1).text = "Years"
    table.cell(1, 0).text = "Python"
    table.cell(1, 1).text = "6"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    """Generates a valid PPTX presentation in-memory."""
    prs = pptx.Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    tx_box = slide.shapes.add_textbox(100, 100, 300, 100)
    tf = tx_box.text_frame
    tf.text = "Slide 1: AI Agent Architecture Overview"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_zip_bytes(sample_pdf_bytes, sample_docx_bytes) -> bytes:
    """Generates a valid ZIP archive containing PDF, DOCX, and TXT files."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("candidate_1.pdf", sample_pdf_bytes)
        z.writestr("candidate_2.docx", sample_docx_bytes)
        z.writestr("notes.txt", "Candidate interview pre-screen notes.")
        # Add a Mac OS artifact that should be filtered out
        z.writestr("__MACOSX/._candidate_1.pdf", b"garbage")
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Unit Tests: Document Parsers
# ---------------------------------------------------------------------------

def test_parse_pdf_valid(sample_pdf_bytes):
    """Test PDF text extraction with PyMuPDF."""
    text = parse_pdf(sample_pdf_bytes)
    assert "Aarav Sharma" in text
    assert "LangGraph" in text
    assert "Groq" in text


def test_parse_pdf_corrupted():
    """Test PDF extraction error handling with corrupted bytes."""
    corrupted_bytes = b"Not a valid PDF byte string"
    text = parse_pdf(corrupted_bytes)
    assert "Error parsing PDF" in text


def test_parse_docx_valid(sample_docx_bytes):
    """Test DOCX extraction including body paragraphs and table cells."""
    text = parse_docx(sample_docx_bytes)
    assert "Lead Systems Architect" in text
    assert "ChromaDB" in text
    assert "Python | 6" in text or "Python" in text


def test_parse_docx_corrupted():
    """Test DOCX parser error handling with non-docx bytes."""
    text = parse_docx(b"bad docx data")
    assert "Error parsing DOCX" in text


def test_parse_pptx_valid(sample_pptx_bytes):
    """Test PPTX extraction from slides."""
    text = parse_pptx(sample_pptx_bytes)
    assert "Slide 1: AI Agent Architecture Overview" in text


def test_parse_zip_valid(sample_zip_bytes):
    """Test ZIP extraction for mixed valid documents while ignoring system artifacts."""
    extracted = parse_zip(sample_zip_bytes)
    filenames = [doc["filename"] for doc in extracted]

    assert len(extracted) == 3
    assert "candidate_1.pdf" in filenames
    assert "candidate_2.docx" in filenames
    assert "notes.txt" in filenames
    assert not any("__MACOSX" in f for f in filenames)

    # Verify extracted text within the zip
    pdf_doc = next(d for d in extracted if d["filename"] == "candidate_1.pdf")
    assert "Aarav Sharma" in pdf_doc["text"]


def test_extract_text_auto(sample_pdf_bytes, sample_docx_bytes):
    """Test auto-routing by file extension."""
    pdf_res = extract_text_auto("resume.pdf", sample_pdf_bytes)
    docx_res = extract_text_auto("resume.docx", sample_docx_bytes)
    txt_res = extract_text_auto("raw.txt", b"Plain text content.")

    assert "Aarav Sharma" in pdf_res
    assert "Lead Systems Architect" in docx_res
    assert txt_res == "Plain text content."


# ---------------------------------------------------------------------------
# Unit Tests: URL & Web Extraction (Mocked Network Calls)
# ---------------------------------------------------------------------------

@patch("trafilatura.fetch_url")
@patch("trafilatura.extract")
def test_parse_url_trafilatura_success(mock_extract, mock_fetch):
    """Test web page extraction via Trafilatura."""
    mock_fetch.return_value = "<html><body>Job Posting</body></html>"
    mock_extract.return_value = "Staff Machine Learning Engineer role requiring 5+ years of Python and PyTorch experience."

    extracted = parse_url("https://careers.example.com/job/123")
    assert "Staff Machine Learning Engineer" in extracted
    mock_fetch.assert_called_once_with("https://careers.example.com/job/123")


@patch("trafilatura.fetch_url", return_value=None)
@patch("requests.get")
def test_parse_url_fallback_beautifulsoup(mock_requests_get, mock_fetch):
    """Test fallback to BeautifulSoup when Trafilatura fails."""
    mock_response = MagicMock()
    mock_response.content = b"""
    <html>
        <head><title>Job Spec</title></head>
        <body>
            <nav>Nav links</nav>
            <h1>Backend AI Engineer</h1>
            <p>We are seeking an engineer proficient in FastAPI and SQLite.</p>
        </body>
    </html>
    """
    mock_requests_get.return_value = mock_response

    extracted = parse_url("https://jobs.example.com/ai-engineer")
    assert "Backend AI Engineer" in extracted
    assert "We are seeking an engineer proficient in FastAPI and SQLite." in extracted
    assert "Nav links" not in extracted  # <nav> tags should be stripped
