import io
import zipfile
import fitz  # PyMuPDF
import docx
import pptx
import trafilatura
from bs4 import BeautifulSoup
import requests
from typing import List, Dict, Any


def parse_pdf(file_bytes: bytes) -> str:
    """Extracts raw textual content from PDF files using PyMuPDF."""
    text_content = []
    try:
        with fitz.open(stream=file_bytes, filetype="pdf") as doc:
            for page in doc:
                text_content.append(page.get_text() or "")
        return "\n".join(text_content).strip()
    except Exception as e:
        return f"Error parsing PDF: {str(e)}"


def parse_docx(file_bytes: bytes) -> str:
    """Extracts text paragraphs and table cells from DOCX files."""
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        # Also parse tables if present
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    table_texts.append(" | ".join(row_text))
                    
        return "\n".join(paragraphs + table_texts).strip()
    except Exception as e:
        return f"Error parsing DOCX: {str(e)}"


def parse_pptx(file_bytes: bytes) -> str:
    """Extracts text runs from presentation slides."""
    try:
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        runs.append(paragraph.text)
        return "\n".join([r for r in runs if r.strip()]).strip()
    except Exception as e:
        return f"Error parsing PPTX: {str(e)}"


def parse_url(url: str) -> str:
    """
    Extracts structured body text from web pages, job boards, or LinkedIn job URLs.
    Uses Trafilatura with a BeautifulSoup fallback.
    """
    try:
        downloaded = trafilatura.fetch_url(url)
        if downloaded:
            extracted = trafilatura.extract(
                downloaded,
                include_comments=False,
                include_tables=True,
                no_fallback=False
            )
            if extracted and len(extracted.strip()) > 50:
                return extracted.strip()

        # Fallback to requests + BeautifulSoup
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        resp = requests.get(url, headers=headers, timeout=10)
        soup = BeautifulSoup(resp.content, "html.parser")
        
        for script_or_style in soup(["script", "style", "nav", "footer", "header"]):
            script_or_style.extract()
            
        paragraphs = [p.get_text() for p in soup.find_all(["p", "h1", "h2", "h3", "li", "span"]) if p.get_text().strip()]
        return "\n".join(paragraphs).strip()
    except Exception as e:
        return f"Error extracting content from URL {url}: {str(e)}"


def parse_zip(zip_bytes: bytes) -> List[Dict[str, str]]:
    """
    Recursively extracts and parses all PDF and DOCX files contained inside a ZIP archive.
    Returns a list of dicts with filenames and extracted text.
    """
    extracted_docs = []
    try:
        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as z:
            for filename in z.namelist():
                if filename.startswith("__MACOSX/") or filename.endswith("/"):
                    continue
                file_data = z.read(filename)
                base_name = filename.split("/")[-1]
                
                if base_name.lower().endswith(".pdf"):
                    text = parse_pdf(file_data)
                    extracted_docs.append({"filename": base_name, "text": text})
                elif base_name.lower().endswith(".docx"):
                    text = parse_docx(file_data)
                    extracted_docs.append({"filename": base_name, "text": text})
                elif base_name.lower().endswith(".txt"):
                    text = file_data.decode("utf-8", errors="ignore")
                    extracted_docs.append({"filename": base_name, "text": text})
        return extracted_docs
    except Exception as e:
        return [{"filename": "error.zip", "text": f"Error unpacking ZIP: {str(e)}"}]


def extract_text_auto(file_name: str, file_bytes: bytes) -> str:
    """Auto-detects file type by extension and delegates parsing."""
    ext = file_name.lower().split(".")[-1]
    if ext == "pdf":
        return parse_pdf(file_bytes)
    elif ext in ["docx", "doc"]:
        return parse_docx(file_bytes)
    elif ext in ["pptx", "ppt"]:
        return parse_pptx(file_bytes)
    elif ext in ["txt", "md"]:
        return file_bytes.decode("utf-8", errors="ignore")
    return ""
